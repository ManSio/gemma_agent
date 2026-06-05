from __future__ import annotations

import logging

import json
import os
import zipfile
from typing import Any, Dict, List

from core.error_analysis import record_error_event

# Расширения, которые parse_file обрабатывает явно или как текст
_KNOWN_DOC_EXTS = frozenset(
    {
        "pdf",
        "docx",
        "doc",
        "pptx",
        "xlsx",
        "zip",
        "log",
        "txt",
        "json",
        "yaml",
        "yml",
        "csv",
        "md",
        "markdown",
        "rst",
        "ini",
        "cfg",
        "toml",
        "env",
        "xml",
        "html",
        "htm",
        "sql",
        "py",
        "js",
        "ts",
        "css",
        "c",
        "h",
        "cpp",
        "go",
        "rs",
        "java",
        "sh",
        "bat",
        "ps1",
    }
)

# Типы с отдельным парсером (не «просто UTF-8 текст»)
_STRUCTURED_DOC_EXTS = frozenset(
    {"pdf", "docx", "doc", "pptx", "xlsx", "zip", "json", "yaml", "yml"}
)
_PLAIN_TEXT_EXTS = _KNOWN_DOC_EXTS - _STRUCTURED_DOC_EXTS


logger = logging.getLogger(__name__)

def _structured_document_body(doc: Dict[str, Any]) -> str:
    """Текст для промпта из XLSX / ZIP / лога / JSON-YAML preview, если нет поля text."""
    t = str(doc.get("text") or "").strip()
    if t:
        return t
    tables = doc.get("tables")
    if isinstance(tables, list) and tables:
        lines: List[str] = []
        for tbl in tables[:15]:
            if not isinstance(tbl, dict):
                continue
            sheet = str(tbl.get("sheet") or "")
            rows = tbl.get("rows") or []
            if sheet:
                lines.append(f"=== {sheet} ===")
            if isinstance(rows, list):
                for row in rows[:80]:
                    if isinstance(row, (list, tuple)):
                        lines.append("\t".join("" if c is None else str(c) for c in row))
                    else:
                        lines.append(str(row))
        return "\n".join(lines)
    entries = doc.get("entries")
    if isinstance(entries, list) and entries:
        return "Архив ZIP, файлы:\n" + "\n".join(str(x) for x in entries[:300])
    if "line_count" in doc and doc.get("ok"):
        parts = [f"Файл-лог: строк ≈ {doc.get('line_count')}"]
        errs = doc.get("errors")
        warns = doc.get("warnings")
        if isinstance(errs, list) and errs:
            parts.append("Строки с error:\n" + "\n".join(str(x) for x in errs[:80]))
        if isinstance(warns, list) and warns:
            parts.append("Строки с warn:\n" + "\n".join(str(x) for x in warns[:80]))
        return "\n".join(parts)
    preview = doc.get("preview")
    if preview is not None and str(preview).strip():
        fmt = doc.get("format") or "data"
        return f"[{fmt}]\n{str(preview).strip()}"
    return ""


def intake_storable_plain(doc: Dict[str, Any]) -> str:
    """Плоский текст/таблицы для сохранения на диск (без префикса [вложение])."""
    if not isinstance(doc, dict) or not doc.get("ok"):
        return ""
    if doc.get("text_layer_empty"):
        return ""
    return _structured_document_body(doc).strip()


def format_document_intake_for_brain(doc: Dict[str, Any], *, max_chars: int | None = None) -> str:
    """
    Текст для user-промпта мозга. Раньше document_intake жил только в context dict, но не попадал в assemble — модель «не видела» файл.
    """
    if not isinstance(doc, dict) or not doc:
        return ""
    if max_chars is None:
        try:
            max_chars = int((os.getenv("BRAIN_DOCUMENT_PROMPT_MAX_CHARS") or "14000").strip())
        except ValueError:
            max_chars = 14000
    max_chars = max(2000, min(max_chars, 120_000))
    if not doc.get("ok"):
        err = str(doc.get("error") or "unknown")
        if err == "worker_timeout_or_failed":
            hint = (
                "файл слишком тяжёлый для разбора за отведённое время. Предложи: меньший PDF, первые страницы отдельным файлом, "
                "или увеличение HEAVY_WORKER_TIMEOUT_SEC в настройках бота."
            )
        elif "timeout" in err.lower():
            hint = "таймаут при разборе; для больших PDF увеличь HEAVY_WORKER_TIMEOUT_SEC или пришли фрагмент."
        elif "python-docx not installed" in err:
            hint = "на сервере не установлен пакет python-docx (зависимость requirements.txt)."
        elif "legacy_doc_format" in err:
            hint = "это старый формат Word (.doc); попроси сохранить как .docx."
        elif "not_valid_docx" in err:
            hint = "файл не похож на корректный .docx; попроси пересохранить из Word/LibreOffice как DOCX."
        elif "python-pptx not installed" in err:
            hint = "на сервере не установлен python-pptx (requirements.txt)."
        elif "openpyxl not installed" in err:
            hint = "на сервере не установлен openpyxl (requirements.txt)."
        elif "pyyaml not installed" in err:
            hint = "на сервере не установлен PyYAML (requirements.txt)."
        else:
            hint = "попроси другой формат или фрагмент текста."
        return (
            "[вложение] Разбор файла не удался. Сообщи пользователю кратко по-русски. "
            f"{hint} Технически: {err}"
        )
    if doc.get("text_layer_empty"):
        return (
            "[вложение] Файл открыт, но извлечённого текста нет (часто PDF-скан или защищённый документ без текстового слоя). "
            "Предложи: текстовый PDF, DOCX, скрин/фото с OCR, или вопрос в подписи к файлу."
        )
    body = _structured_document_body(doc).strip()
    if not body:
        return (
            "[вложение] Текст из файла пустой. Возможен скан или неподдерживаемая структура. "
            "Предложи прислать текстовый вариант или уточнить задачу в сообщении."
        )
    head = f"[вложение] Извлечённый текст (до {max_chars} символов):\n"
    if len(body) <= max_chars:
        return head + body
    return head + body[:max_chars] + "\n… (обрезано; увеличь BRAIN_DOCUMENT_PROMPT_MAX_CHARS или попроси фрагмент)"


class DocumentIntakeModule:
    def __init__(self) -> None:
        self.enabled = os.getenv("DOC_INTAKE_ENABLED", "true").strip().lower() in {"1", "true", "yes", "on"}
        self.max_zip_entries = int(os.getenv("ZIP_MAX_ENTRIES", "200"))
        self.max_zip_unpacked_mb = float(os.getenv("ZIP_MAX_UNPACKED_MB", "100"))
        self.pdf_max_pages = max(1, int(os.getenv("PDF_INTAKE_MAX_PAGES", "200")))
        # auto | pymupdf | pypdf — pymupdf быстрее на крупных PDF и реже упирается в таймаут воркера
        _b = (os.getenv("PDF_INTAKE_BACKEND") or "auto").strip().lower()
        self._pdf_backend = _b if _b in {"auto", "pymupdf", "pypdf"} else "auto"

    def _read_text_file(self, path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def _pdf_text_pymupdf(self, path: str) -> str | None:
        """Возвращает извлечённый текст или None, если бэкенд недоступен/ошибка (тогда пробуем pypdf)."""
        try:
            import fitz  # PyMuPDF  # type: ignore
        except Exception:
            return None
        doc = None
        try:
            doc = fitz.open(path)
            if getattr(doc, "needs_pass", False):
                return None
            n = min(int(doc.page_count), self.pdf_max_pages)
            parts: List[str] = []
            for i in range(n):
                page = doc.load_page(i)
                parts.append(page.get_text() or "")
            return "\n".join(parts)
        except Exception as e:
            record_error_event("document_intake", "pymupdf extract failed", exc=e, extra={"path": path})
            return None
        finally:
            if doc is not None:
                try:
                    doc.close()
                except Exception as e:
                    logger.debug('%s optional failed: %s', 'document_intake', e, exc_info=True)
    def _pdf_text_pypdf(self, path: str) -> str:
        import pypdf  # type: ignore

        with open(path, "rb") as f:
            reader = pypdf.PdfReader(f)
            return "\n".join((p.extract_text() or "") for p in reader.pages[: self.pdf_max_pages])

    def parse_pdf(self, path: str) -> Dict[str, Any]:
        text = ""
        backend = self._pdf_backend
        if backend in {"auto", "pymupdf"}:
            raw = self._pdf_text_pymupdf(path)
            if raw is not None:
                text = (raw or "").strip()
            if backend == "pymupdf":
                if not text:
                    return {"ok": True, "text": "", "text_layer_empty": True}
                return {"ok": True, "text": text[:200000]}
        if not text and backend in {"auto", "pypdf"}:
            try:
                raw2 = self._pdf_text_pypdf(path)
                text = (raw2 or "").strip()
            except Exception as e:
                record_error_event("document_intake", "pdf parse failed", exc=e, extra={"path": path})
                return {"ok": False, "error": str(e)}
        if not text:
            return {"ok": True, "text": "", "text_layer_empty": True}
        return {"ok": True, "text": text[:200000]}

    def _docx_collect_text(self, doc: Any) -> str:
        """Абзацы, таблицы (часто вся суть претензий/форм в ячейках), колонтитулы."""
        chunks: List[str] = []
        for p in doc.paragraphs:
            t = (getattr(p, "text", None) or "").strip()
            if t:
                chunks.append(t)
        for table in doc.tables:
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    ct = " ".join(
                        (pp.text or "").strip() for pp in cell.paragraphs if (pp.text or "").strip()
                    )
                    if ct:
                        cells.append(ct)
                if cells:
                    chunks.append(" | ".join(cells))
        for sec in doc.sections:
            for hdr in (sec.header, sec.footer):
                try:
                    for p in hdr.paragraphs:
                        t = (p.text or "").strip()
                        if t:
                            chunks.append(t)
                    for table in hdr.tables:
                        for row in table.rows:
                            cells = []
                            for cell in row.cells:
                                ct = " ".join(
                                    (pp.text or "").strip()
                                    for pp in cell.paragraphs
                                    if (pp.text or "").strip()
                                )
                                if ct:
                                    cells.append(ct)
                            if cells:
                                chunks.append(" | ".join(cells))
                except Exception as e:
                    logger.debug('%s optional failed: %s', 'document_intake', e, exc_info=True)
        return "\n".join(chunks)

    def parse_docx(self, path: str) -> Dict[str, Any]:
        try:
            import docx  # type: ignore
        except ImportError as e:
            record_error_event("document_intake", "python-docx missing", exc=e, extra={"path": path})
            return {
                "ok": False,
                "error": "python-docx not installed (pip install python-docx)",
            }
        try:
            doc = docx.Document(path)
            text = self._docx_collect_text(doc)
            text = (text or "").strip()
            if not text:
                return {"ok": True, "text": "", "text_layer_empty": True}
            return {"ok": True, "text": text[:200000]}
        except Exception as e:
            record_error_event("document_intake", "docx parse failed", exc=e, extra={"path": path})
            err = str(e)
            if "not a zip file" in err.lower() or "bad zipfile" in err.lower():
                return {
                    "ok": False,
                    "error": "not_valid_docx_old_doc_or_corrupt_save_as_docx",
                }
            return {"ok": False, "error": err}

    def parse_pptx(self, path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as e:
            record_error_event("document_intake", "python-pptx missing", exc=e, extra={"path": path})
            return {"ok": False, "error": "python-pptx not installed (pip install python-pptx)"}
        try:
            prs = Presentation(path)
            chunks: List[str] = []
            for slide in prs.slides:
                for shp in slide.shapes:
                    txt = getattr(shp, "text", "")
                    if txt:
                        chunks.append(txt)
            text = "\n".join(chunks).strip()
            if not text:
                return {"ok": True, "text": "", "text_layer_empty": True}
            return {"ok": True, "text": text[:200000]}
        except Exception as e:
            record_error_event("document_intake", "pptx parse failed", exc=e, extra={"path": path})
            err = str(e)
            if "not a zip file" in err.lower() or "bad zipfile" in err.lower():
                return {"ok": False, "error": "not_valid_pptx_corrupt"}
            return {"ok": False, "error": err}

    def parse_xlsx(self, path: str) -> Dict[str, Any]:
        try:
            import openpyxl  # type: ignore
        except ImportError as e:
            record_error_event("document_intake", "openpyxl missing", exc=e, extra={"path": path})
            return {"ok": False, "error": "openpyxl not installed (pip install openpyxl)"}
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            tables = []
            for ws in wb.worksheets[:10]:
                rows = []
                for row in ws.iter_rows(min_row=1, max_row=100, values_only=True):
                    rows.append(list(row))
                tables.append({"sheet": ws.title, "rows": rows})
            try:
                wb.close()
            except Exception as e:
                logger.debug('%s optional failed: %s', 'document_intake', e, exc_info=True)
            return {"ok": True, "tables": tables}
        except Exception as e:
            record_error_event("document_intake", "xlsx parse failed", exc=e, extra={"path": path})
            return {"ok": False, "error": str(e)}

    def parse_zip_safe(self, path: str) -> Dict[str, Any]:
        try:
            total_unpacked = 0
            names: List[str] = []
            with zipfile.ZipFile(path, "r") as zf:
                infos = zf.infolist()
                if len(infos) > self.max_zip_entries:
                    return {"ok": False, "error": "too_many_entries"}
                for info in infos:
                    if ".." in info.filename.replace("\\", "/"):
                        return {"ok": False, "error": "zip_path_traversal_detected"}
                    total_unpacked += int(info.file_size)
                    if total_unpacked > int(self.max_zip_unpacked_mb * 1024 * 1024):
                        return {"ok": False, "error": "zip_unpacked_size_limit"}
                    names.append(info.filename)
            return {"ok": True, "entries": names, "total_unpacked_bytes": total_unpacked}
        except Exception as e:
            record_error_event("document_intake", "zip parse failed", exc=e, extra={"path": path})
            return {"ok": False, "error": str(e)}

    def parse_json_yaml(self, path: str) -> Dict[str, Any]:
        ext = (path.rsplit(".", 1)[-1].lower() if "." in path else "")
        try:
            raw = self._read_text_file(path)
            if ext == "json":
                obj = json.loads(raw)
                return {"ok": True, "format": "json", "valid": True, "preview": str(obj)[:20000]}
            if ext in {"yaml", "yml"}:
                try:
                    import yaml  # type: ignore
                except ImportError as e:
                    record_error_event("document_intake", "pyyaml missing", exc=e, extra={"path": path})
                    return {"ok": False, "error": "pyyaml not installed (pip install PyYAML)"}
                obj = yaml.safe_load(raw)
                return {"ok": True, "format": "yaml", "valid": True, "preview": str(obj)[:20000]}
            return {"ok": False, "error": "unsupported_format"}
        except Exception as e:
            record_error_event("document_intake", "json/yaml parse failed", exc=e, extra={"path": path})
            return {"ok": False, "error": str(e)}

    def parse_plain_text(self, path: str) -> Dict[str, Any]:
        try:
            text = self._read_text_file(path)
            text = (text or "").strip()
            if not text:
                return {"ok": True, "text": "", "text_layer_empty": True}
            return {"ok": True, "text": text[:200000]}
        except Exception as e:
            record_error_event("document_intake", "plain text read failed", exc=e, extra={"path": path})
            return {"ok": False, "error": str(e)}

    def parse_html(self, path: str) -> Dict[str, Any]:
        """HTML/MHTML-сохранения: вытаскиваем текст, а не сырой разметку в промпт."""
        from core.url_fetch import _body_to_text

        try:
            max_b = min(2_000_000, max(4096, int(os.getenv("URL_FETCH_MAX_BYTES", "2097152"))))
            with open(path, "rb") as f:
                raw = f.read(max_b)
        except OSError as e:
            record_error_event("document_intake", "html read failed", exc=e, extra={"path": path})
            return {"ok": False, "error": str(e)}
        text = _body_to_text(raw, "text/html; charset=utf-8").strip()
        if not text:
            return {"ok": True, "text": "", "text_layer_empty": True}
        return {"ok": True, "text": text[:200000]}

    def _sniff_office_zip_ext(self, path: str) -> str:
        try:
            with zipfile.ZipFile(path, "r") as zf:
                names = {n.replace("\\", "/") for n in zf.namelist()}
        except (zipfile.BadZipFile, OSError, ValueError):
            return ""
        if "word/document.xml" in names or any(
            n.startswith("word/") and "document" in n for n in names
        ):
            return "docx"
        if "[Content_Types].xml" in names and any("word/" in n for n in names):
            return "docx"
        if "xl/workbook.xml" in names:
            return "xlsx"
        if "ppt/presentation.xml" in names or any(n.startswith("ppt/slides/") for n in names):
            return "pptx"
        return ""

    def _try_plaintext_binary(self, path: str) -> Dict[str, Any] | None:
        max_b = min(400_000, int(os.getenv("DOC_INTAKE_PLAIN_MAX_BYTES", "400000")))
        try:
            with open(path, "rb") as f:
                raw = f.read(max_b)
        except OSError:
            return None
        if not raw:
            return {"ok": True, "text": "", "text_layer_empty": True}
        text: str | None = None
        for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
            try:
                text = raw.decode(enc)
                break
            except UnicodeDecodeError:
                continue
        if text is None:
            return None
        sample = text[: min(8000, len(text))]
        if len(sample) > 12:
            printable = sum(1 for c in sample if c.isprintable() or c in "\n\r\t\u00a0")
            if printable / len(sample) < 0.88:
                return None
        text = text.strip()
        if not text:
            return {"ok": True, "text": "", "text_layer_empty": True}
        return {"ok": True, "text": text[:200000]}

    def parse_file(self, path: str) -> Dict[str, Any]:
        ext = (path.rsplit(".", 1)[-1].lower() if "." in path else "")
        if ext not in _KNOWN_DOC_EXTS:
            sniffed = self._sniff_office_zip_ext(path)
            if sniffed:
                ext = sniffed
        # Старый Word .doc (OLE), не OOXML — python-docx не откроет
        if ext == "doc":
            try:
                with open(path, "rb") as f:
                    sig = f.read(8)
                if sig.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
                    return {
                        "ok": False,
                        "error": "legacy_doc_format_save_as_docx",
                    }
            except OSError:
                pass
        if ext == "pdf":
            return self.parse_pdf(path)
        if ext == "docx":
            return self.parse_docx(path)
        if ext == "pptx":
            return self.parse_pptx(path)
        if ext == "xlsx":
            return self.parse_xlsx(path)
        if ext == "zip":
            return self.parse_zip_safe(path)
        if ext in {"json", "yaml", "yml"}:
            return self.parse_json_yaml(path)
        if ext in {"html", "htm"}:
            return self.parse_html(path)
        if ext in _PLAIN_TEXT_EXTS:
            try:
                with open(path, "rb") as f:
                    head = f.read(8000)
                low = head.lstrip().lower()
                if low.startswith(b"<!doctype html") or low.startswith(b"<html"):
                    return self.parse_html(path)
            except OSError:
                pass
            return self.parse_plain_text(path)
        sniffed2 = self._sniff_office_zip_ext(path)
        if sniffed2:
            if sniffed2 == "docx":
                return self.parse_docx(path)
            if sniffed2 == "xlsx":
                return self.parse_xlsx(path)
            if sniffed2 == "pptx":
                return self.parse_pptx(path)
        pt = self._try_plaintext_binary(path)
        if pt is not None:
            return pt
        try:
            with open(path, "rb") as f:
                head = f.read(12000)
            low = head.lstrip().lower()
            if low.startswith(b"<!doctype html") or low.startswith(b"<html") or b"<html" in low[:4000]:
                return self.parse_html(path)
        except OSError:
            pass
        return {"ok": False, "error": f"unsupported_extension:{ext}"}
