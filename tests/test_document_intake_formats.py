import os
import tempfile

from core.document_intake import DocumentIntakeModule, format_document_intake_for_brain


def test_xlsx_tables_surface_in_brain_prompt():
    openpyxl = __import__("openpyxl")
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "t.xlsx")
        wb = openpyxl.Workbook()
        ws = wb.active
        assert ws is not None
        ws.title = "S1"
        ws.append(["a", "b"])
        ws.append([1, 2])
        wb.save(path)
        wb.close()

        mod = DocumentIntakeModule()
        doc = mod.parse_file(path)
        assert doc.get("ok") is True
        s = format_document_intake_for_brain(doc, max_chars=5000)
        assert "S1" in s
        assert "1" in s


def test_plain_unknown_extension():
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "readme")  # без расширения
        with open(path, "w", encoding="utf-8") as f:
            f.write("hello plain no ext\n")
        mod = DocumentIntakeModule()
        doc = mod.parse_file(path)
        assert doc.get("ok") is True
        assert "hello plain" in (doc.get("text") or "")


def test_html_strips_markup():
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "p.html")
        with open(path, "w", encoding="utf-8") as f:
            f.write(
                "<!DOCTYPE html><html><head><title>T</title></head>"
                "<body><p>Visible line</p><script>no</script></body></html>"
            )
        mod = DocumentIntakeModule()
        doc = mod.parse_file(path)
        assert doc.get("ok") is True
        text = doc.get("text") or ""
        assert "Visible line" in text
        assert "<p>" not in text


def test_sniff_html_file_without_extension():
    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "saved_page")
        with open(path, "w", encoding="utf-8") as f:
            f.write("<html><body><h1>Headline</h1></body></html>")
        mod = DocumentIntakeModule()
        doc = mod.parse_file(path)
        assert doc.get("ok") is True
        assert "Headline" in (doc.get("text") or "")


def test_pptx_roundtrip():
    pptx = __import__("pptx")
    from pptx.util import Inches  # type: ignore

    with tempfile.TemporaryDirectory() as td:
        path = os.path.join(td, "p.pptx")
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        box.text = "Слайд тест"
        prs.save(path)

        mod = DocumentIntakeModule()
        doc = mod.parse_file(path)
        assert doc.get("ok") is True
        assert "Слайд" in (doc.get("text") or "")
